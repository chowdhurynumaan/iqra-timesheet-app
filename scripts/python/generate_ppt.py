#!/usr/bin/env python3
"""
Generate PowerPoint slide for TimeSheet Dashboard User Guide
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_user_guide_slide():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "📊 TimeSheet Attendance Dashboard - Complete User Guide"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(100, 126, 234)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.75), Inches(9.4), Inches(0.35))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Upload reports → Detect month → Review data → Save to cloud → Edit as needed | All steps tracked with full audit history"
    subtitle_p.font.size = Pt(12)
    subtitle_p.font.color.rgb = RGBColor(102, 102, 102)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Left column - HOW TO USE
    left_box = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2), Inches(3.1), Inches(5.9))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(249, 249, 249)
    left_box.line.color.rgb = RGBColor(102, 126, 234)
    left_box.line.width = Pt(2)
    
    how_to_title = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(2.9), Inches(0.4))
    how_to_frame = how_to_title.text_frame
    how_to_p = how_to_frame.paragraphs[0]
    how_to_p.text = "🚀 HOW TO USE (5 Steps)"
    how_to_p.font.size = Pt(14)
    how_to_p.font.bold = True
    how_to_p.font.color.rgb = RGBColor(102, 126, 234)
    
    steps_text = """1) Open: https://chowdhurynumaan.github.io/iqra-timesheet-app/

2) Select Month: Choose from dropdown (auto-populates)

3) Upload Report: Click Upload → Select Excel/CSV file

4) Review Data: Check employee cards, metrics, compare

5) Save: Click Save to Cloud → Confirm changes

✏️ BONUS: Click Edit on any card to modify data

💡 Tip: Month auto-detected. Data in cloud forever. File not needed after save."""
    
    steps_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(2.9), Inches(5.0))
    steps_frame = steps_box.text_frame
    steps_frame.word_wrap = True
    steps_frame.text = steps_text
    for paragraph in steps_frame.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
        paragraph.space_before = Pt(3)
        paragraph.space_after = Pt(3)
    
    # Middle column - ERRORS
    middle_box = slide.shapes.add_shape(1, Inches(3.5), Inches(1.2), Inches(3.1), Inches(5.9))
    middle_box.fill.solid()
    middle_box.fill.fore_color.rgb = RGBColor(255, 243, 205)
    middle_box.line.color.rgb = RGBColor(255, 107, 107)
    middle_box.line.width = Pt(2)
    
    errors_title = slide.shapes.add_textbox(Inches(3.6), Inches(1.3), Inches(2.9), Inches(0.4))
    errors_frame = errors_title.text_frame
    errors_p = errors_frame.paragraphs[0]
    errors_p.text = "⚠️ COMMON ERRORS"
    errors_p.font.size = Pt(14)
    errors_p.font.bold = True
    errors_p.font.color.rgb = RGBColor(255, 107, 107)
    
    errors_text = """❌ 404 Error: Hard refresh (Ctrl+F5). Site rebuilds in ~1 min.

❌ Upload fails: Check file is .xls/.xlsx/.csv with headers

❌ No month detected: Add "DD-MM-YYYY ~ DD-MM-YYYY" format

❌ Data missing cols: Ensure columns: Date, Name, Hours, Late

❌ Save doesn't work: Check internet. Open Console (F12)

❌ Firebase error: Clear cache. Check network tab.

✅ Still stuck? Open browser console (F12 key) for detailed error messages."""
    
    errors_box = slide.shapes.add_textbox(Inches(3.6), Inches(1.8), Inches(2.9), Inches(5.0))
    errors_box_frame = errors_box.text_frame
    errors_box_frame.word_wrap = True
    errors_box_frame.text = errors_text
    for paragraph in errors_box_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
        paragraph.space_before = Pt(2)
        paragraph.space_after = Pt(2)
    
    # Right column - FEATURES
    right_box = slide.shapes.add_shape(1, Inches(6.7), Inches(1.2), Inches(3.1), Inches(5.9))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(227, 242, 253)
    right_box.line.color.rgb = RGBColor(33, 150, 243)
    right_box.line.width = Pt(2)
    
    features_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(2.9), Inches(0.4))
    features_frame = features_title.text_frame
    features_p = features_frame.paragraphs[0]
    features_p.text = "✨ KEY FEATURES"
    features_p.font.size = Pt(14)
    features_p.font.bold = True
    features_p.font.color.rgb = RGBColor(33, 150, 243)
    
    features_text = """📊 Analytics: Attendance %, late mins, dept summaries

📅 Smart Detection: Auto-detects month from file

🔄 Merge Logic: Shows NEW/MODIFIED/UNCHANGED

📋 Daily Breakdown: Punch times per employee

📈 Charts: Visual trends (Chart.js powered)

🖨️ Export: CSV, Excel, or print format

🔐 Secure Cloud: Encrypted in Firebase

📝 Audit Trail: All edits logged with timestamps

🧹 Cleanup: Remove duplicates with one click

⚙️ Settings: Customize work hours per month"""
    
    features_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(2.9), Inches(5.0))
    features_box_frame = features_box.text_frame
    features_box_frame.word_wrap = True
    features_box_frame.text = features_text
    for paragraph in features_box_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
        paragraph.space_before = Pt(2)
        paragraph.space_after = Pt(2)
    
    # Bottom footer
    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(9.4), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "📌 Remember: All edits tracked | Cloud keeps history | Old files safe in GitHub | Settings per month | Questions? Open browser console (F12) for detailed errors"
    footer_p.font.size = Pt(9)
    footer_p.font.color.rgb = RGBColor(102, 102, 102)
    footer_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "TimeSheet_User_Guide.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_user_guide_slide()
