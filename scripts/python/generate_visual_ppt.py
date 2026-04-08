#!/usr/bin/env python3
"""
Generate visual PowerPoint with diagrams, flowcharts, and mockups
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import io
import tempfile
import os

def create_workflow_diagram():
    """Create workflow diagram image"""
    img = Image.new('RGB', (1200, 250), 'white')
    draw = ImageDraw.Draw(img)
    
    # Colors
    color_step = RGBColor(102, 126, 234)
    color_arrow = RGBColor(150, 150, 150)
    
    steps = [
        "📁 Upload\nExcel File",
        "🔍 Auto-Detect\nMonth",
        "📊 Review\nData",
        "☁️ Save to\nCloud",
        "✏️ Edit &\nTrack"
    ]
    
    # Draw steps
    x_start = 50
    x_gap = 230
    y_center = 125
    box_width = 150
    box_height = 200
    
    for i, step in enumerate(steps):
        x = x_start + (i * x_gap)
        
        # Draw box
        draw.rectangle(
            [x, y_center - 50, x + box_width, y_center + 50],
            fill='#f0f2f5',
            outline='#667eea',
            width=3
        )
        
        # Draw text
        draw.text(
            (x + box_width//2, y_center),
            step,
            fill='#333',
            anchor='mm',
            align='center'
        )
        
        # Draw arrow
        if i < len(steps) - 1:
            arrow_x = x + box_width + 15
            draw.polygon(
                [(arrow_x, y_center), (arrow_x + 25, y_center - 8), (arrow_x + 25, y_center + 8)],
                fill='#667eea'
            )
            draw.line(
                [(x + box_width, y_center), (arrow_x, y_center)],
                fill='#667eea',
                width=2
            )
    
    return img

def create_error_mockup():
    """Create error screen mockup"""
    img = Image.new('RGB', (500, 600), '#f0f2f5')
    draw = ImageDraw.Draw(img)
    
    # Browser chrome
    draw.rectangle([0, 0, 500, 40], fill='#333')
    draw.text((10, 8), 'chowdhurynumaan.github.io/iqra-timesheet-app/', fill='#666', anchor='lm')
    
    # 404 Error
    draw.rectangle([20, 80, 480, 250], fill='white', outline='#ff6b6b', width=2)
    draw.text((250, 110), '404', fill='#ff6b6b', anchor='mm')
    draw.text((250, 150), 'File not found', fill='#666', anchor='mm')
    draw.text((250, 190), 'reports/dashboard.html not found', fill='#999', anchor='mm')
    
    # Solution
    draw.rectangle([20, 280, 480, 550], fill='#e8f5e9', outline='#28a745', width=2)
    draw.text((30, 310), '✓ SOLUTION:', fill='#28a745', anchor='lm')
    draw.text((30, 350), '1. Hard refresh (Ctrl + F5)', fill='#333', anchor='lm')
    draw.text((30, 380), '2. Clear browser cache', fill='#333', anchor='lm')
    draw.text((30, 410), '3. Wait 1 minute for rebuild', fill='#333', anchor='lm')
    draw.text((30, 450), '4. File moved to src/dashboard.html', fill='#333', anchor='lm')
    draw.text((30, 490), '→ Auto-redirects now', fill='#28a745', anchor='lm')
    
    return img

def create_dashboard_mockup():
    """Create dashboard interface mockup"""
    img = Image.new('RGB', (600, 700), 'white')
    draw = ImageDraw.Draw(img)
    
    # Header
    draw.rectangle([0, 0, 600, 80], fill='#667eea')
    draw.text((300, 20), '📊 Attendance Dashboard', fill='white', anchor='mm')
    draw.text((300, 50), 'TimeSheet Portal', fill='#ccc', anchor='mm')
    
    # Controls
    y = 110
    draw.rectangle([20, y, 280, y+40], fill='#f0f2f5', outline='#667eea', width=2)
    draw.text((30, y+20), 'Select Month: Mar 2026', fill='#333', anchor='lm')
    
    draw.rectangle([300, y, 560, y+40], fill='#667eea', outline='#667eea', width=2)
    draw.text((430, y+20), '📄 Upload Report', fill='white', anchor='mm')
    
    # Employee cards
    y = 170
    for i in range(2):
        card_y = y + (i * 150)
        draw.rectangle([20, card_y, 560, card_y+130], fill='#f9f9f9', outline='#667eea', width=2)
        draw.text((30, card_y+15), f'👤 Employee {i+1}', fill='#333', anchor='lm')
        draw.text((30, card_y+45), 'Hours: 160  |  Late: 45 min  |  Absence: 1 day', fill='#666', anchor='lm')
        draw.text((30, card_y+75), 'Dept: Sales  |  Status: ✓ Present', fill='#28a745', anchor='lm')
        draw.rectangle([450, card_y+40, 550, card_y+70], fill='#667eea', outline='#667eea', width=1)
        draw.text((500, card_y+55), '✏️ Edit', fill='white', anchor='mm')
    
    # Footer buttons
    y = 520
    draw.rectangle([20, y, 280, y+40], fill='#28a745', outline='#28a745', width=2)
    draw.text((150, y+20), '☁️ Save to Cloud', fill='white', anchor='mm')
    
    draw.rectangle([300, y, 560, y+40], fill='#2196F3', outline='#2196F3', width=2)
    draw.text((430, y+20), '📊 Analytics', fill='white', anchor='mm')
    
    return img

def create_features_visual():
    """Create features visual grid"""
    img = Image.new('RGB', (600, 400), 'white')
    draw = ImageDraw.Draw(img)
    
    features = [
        ('📊', 'Analytics', 'Dept summaries\n& trends'),
        ('🔄', 'Auto-Merge', 'NEW/MODIFIED\ndetection'),
        ('📝', 'Audit Trail', 'All edits\nlogged'),
        ('🔐', 'Secure', 'Firebase\nencrypted'),
        ('📈', 'Charts', 'Visual\ndata'),
        ('🖨️', 'Export', 'CSV, XLS\nprint'),
    ]
    
    for idx, (emoji, title, desc) in enumerate(features):
        row = idx // 3
        col = idx % 3
        x = 20 + (col * 190)
        y = 20 + (row * 180)
        
        # Card
        draw.rectangle([x, y, x+170, y+160], fill='#f0f2f5', outline='#667eea', width=2)
        
        # Emoji
        draw.text((x+85, y+20), emoji, fill='#667eea', anchor='mm')
        
        # Title
        draw.text((x+85, y+50), title, fill='#333', anchor='mm')
        
        # Description
        draw.text((x+85, y+100), desc, fill='#666', anchor='mm')
    
    return img

def create_visual_ppt():
    """Create visual PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "📊 TimeSheet Dashboard - Visual User Guide"
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(102, 126, 234)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Create temp directory for images
    temp_dir = tempfile.gettempdir()
    
    # === WORKFLOW DIAGRAM ===
    wf_img = create_workflow_diagram()
    wf_path = os.path.join(temp_dir, 'workflow.png')
    wf_img.save(wf_path, format='PNG')
    slide.shapes.add_picture(wf_path, Inches(0.3), Inches(0.65), width=Inches(9.4))
    
    # Workflow label
    wf_label = slide.shapes.add_textbox(Inches(0.3), Inches(1.8), Inches(9.4), Inches(0.3))
    wf_label_frame = wf_label.text_frame
    wf_label_p = wf_label_frame.paragraphs[0]
    wf_label_p.text = "⬆️ 5-STEP WORKFLOW - Follow arrows from left to right"
    wf_label_p.font.size = Pt(12)
    wf_label_p.font.bold = True
    wf_label_p.font.color.rgb = RGBColor(102, 126, 234)
    wf_label_p.alignment = PP_ALIGN.CENTER
    
    # === LEFT COLUMN: ERROR HANDLING ===
    error_img = create_error_mockup()
    error_path = os.path.join(temp_dir, 'error.png')
    error_img.save(error_path, format='PNG')
    slide.shapes.add_picture(error_path, Inches(0.3), Inches(2.15), width=Inches(3.1), height=Inches(4.7))
    
    # === CENTER COLUMN: DASHBOARD MOCKUP ===
    dash_img = create_dashboard_mockup()
    dash_path = os.path.join(temp_dir, 'dashboard.png')
    dash_img.save(dash_path, format='PNG')
    slide.shapes.add_picture(dash_path, Inches(3.45), Inches(2.15), width=Inches(3.1), height=Inches(4.7))
    
    # === RIGHT COLUMN: FEATURES ===
    feat_img = create_features_visual()
    feat_path = os.path.join(temp_dir, 'features.png')
    feat_img.save(feat_path, format='PNG')
    slide.shapes.add_picture(feat_path, Inches(6.6), Inches(2.15), width=Inches(3.1), height=Inches(4.7))
    
    # Cleanup temp files
    for path in [wf_path, error_path, dash_path, feat_path]:
        try:
            os.remove(path)
        except:
            pass
    
    # Bottom footer
    footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.35))
    footer_frame = footer.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "🔗 https://chowdhurynumaan.github.io/iqra-timesheet-app/  |  💾 All data saved to cloud with audit trail  |  ✅ Works: Chrome, Firefox, Safari, Edge"
    footer_p.font.size = Pt(10)
    footer_p.font.color.rgb = RGBColor(102, 102, 102)
    footer_p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = "TimeSheet_User_Guide_Visual.pptx"
    prs.save(output_path)
    print(f"✅ Visual PowerPoint created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_visual_ppt()
